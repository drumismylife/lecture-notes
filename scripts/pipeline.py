#!/usr/bin/env python3
"""
pipeline.py — 강의노트 자동화 파이프라인 (4단계)

사용법:
  python scripts/pipeline.py

input/ 폴더 구조:
  input/[과목명]/weekNN/
    transcript.txt   ← 녹취스크립트 (필수)
    slides.pdf        ← 강의자료 (선택)
    slides.pptx       ← PPT (선택)

처리 완료된 폴더는 input/done/[과목]/weekNN_YYYYMMDD/ 로 보관됩니다.
"""

import os
import re
import sys
import shutil
import subprocess
from pathlib import Path
from datetime import datetime

import anthropic

try:
    from pdfminer.high_level import extract_text as _pdf_extract
    HAS_PDFMINER = True
except ImportError:
    HAS_PDFMINER = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ── 경로 ─────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent
ROOT_DIR   = SCRIPT_DIR.parent
INPUT_DIR  = ROOT_DIR / "input"
DONE_DIR   = INPUT_DIR / "done"

sys.path.insert(0, str(SCRIPT_DIR))
from update_data import update as _update_data_js, COURSE_MAP
from merge import CSS

# 출력 폴더가 과목명과 다른 경우
OUTPUT_FOLDER_OVERRIDE = {"헬라어": "greek"}

# ── CSS 포함 시스템 프롬프트 ──────────────────────────────────
SYSTEM_PROMPT = f"""당신은 신학대학원 강의노트를 고품질 HTML로 정리하는 전문가입니다.

## 출력 형식
반드시 아래 CSS와 HTML 구조를 그대로 사용하세요.
완전한 HTML 파일만 출력하세요 (```html 코드블록 없이 <!DOCTYPE html>로 시작).

CSS:
<style>
{CSS}
</style>

HTML 구조 (필수 준수):
```
<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>{{과목명}} — {{주차}}</title>
<style>/* 위 CSS 그대로 */</style>
</head>
<body>
<div class="page">
  <div class="header">
    <div class="week-badge">{{주차}} · {{날짜}}</div>
    <div class="course-label">한국침례신학대학교 · {{과목명}}</div>
    <h1>{{강의제목 앞부분}} <span>{{강의제목 마지막 단어}}</span></h1>
    <div class="header-meta">{{교수명 · 시간 (있으면)}}</div>
  </div>
  <div class="ornament">✦ ✦ ✦</div>

  <!-- 섹션들 (4~6개 권장) -->
  <div class="section">
    <div class="section-header"><span class="section-icon">📚</span><h2>강의 핵심 정보</h2></div>
    <div class="section-body"><table>...</table></div>
  </div>
  <div class="section">
    <div class="section-header"><span class="section-icon">📖</span><h2>강의 내용</h2></div>
    <div class="section-body"><div class="prose">...</div></div>
  </div>
  <div class="section">
    <div class="section-header"><span class="section-icon">🔑</span><h2>핵심 포인트</h2></div>
    <div class="section-body">
      <div class="key-points">
        <div class="key-point">
          <div class="key-num">01</div>
          <div class="key-text"><strong>제목</strong><p>설명</p></div>
        </div>
      </div>
    </div>
  </div>

  <!-- 다음 주 정보 있으면 추가 -->
  <div class="next-week">
    <div class="label">Next Week · {{다음주 날짜}}</div>
    <h3>{{다음 주제}}</h3>
    <p>{{준비사항}}</p>
  </div>
</div>
</body>
</html>
```

## 사용 가능한 컴포넌트
- `.prose > p` : 일반 단락
- `table > thead + tbody` : 표 (첫 행이 항목/구분/내용이면 thead로)
- `.key-points > .key-point` : 번호형 핵심 포인트 (key-num + key-text)
- `.quote-block > p` + `.quote-sub` : 인용문
- `.callout` + `.callout-title` : 참고/팁 박스
- `.alert` + `.alert-icon` : 경고/주의 박스
- `.note-list > li` : 화살표 목록
- `.keyword-list > .keyword-tag` : 키워드 태그

## 작성 규칙
1. 녹취스크립트를 중심으로 핵심 내용만 추출·정리
2. 강의 내용을 논리적으로 섹션 분리 (보통 4~6개)
3. 표는 항목이 2개 이상일 때 사용
4. 다음 주 정보가 없으면 .next-week 생략
5. 한국어로 작성 (영어 학술용어 그대로 사용 가능)
"""


# ── 텍스트 추출 ───────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    if not HAS_PDFMINER:
        print(f"    ⚠️  pdfminer.six 없음 — {path.name} 스킵 (pip install pdfminer.six)")
        return ""
    try:
        return _pdf_extract(str(path)) or ""
    except Exception as e:
        print(f"    ⚠️  PDF 추출 실패 ({path.name}): {e}")
        return ""


def _extract_pptx(path: Path) -> str:
    if not HAS_PPTX:
        print(f"    ⚠️  python-pptx 없음 — {path.name} 스킵 (pip install python-pptx)")
        return ""
    try:
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        print(f"    ⚠️  PPTX 추출 실패 ({path.name}): {e}")
        return ""


def collect_inputs(week_dir: Path) -> dict:
    transcript_parts, pdf_parts, pptx_parts = [], [], []

    for f in sorted(week_dir.iterdir()):
        ext = f.suffix.lower()
        if ext in (".txt", ".md"):
            transcript_parts.append(f.read_text(encoding="utf-8", errors="ignore"))
        elif ext == ".pdf":
            text = _extract_pdf(f)
            if text.strip():
                pdf_parts.append(f"[{f.name}]\n{text.strip()}")
        elif ext in (".pptx", ".ppt"):
            text = _extract_pptx(f)
            if text.strip():
                pptx_parts.append(f"[{f.name}]\n{text.strip()}")

    return {
        "transcript": "\n\n".join(transcript_parts),
        "pdf": pdf_parts,
        "pptx": pptx_parts,
    }


# ── Claude API 호출 ──────────────────────────────────────────

def generate_html(subject: str, week_num: int, inputs: dict) -> str:
    client = anthropic.Anthropic()

    parts = [f"과목: {subject}\n주차: Week {week_num:02d}\n"]

    if inputs["transcript"].strip():
        parts.append(f"\n## 녹취스크립트\n\n{inputs['transcript']}")

    for text in inputs["pdf"]:
        parts.append(f"\n## PDF 강의자료\n\n{text}")

    for text in inputs["pptx"]:
        parts.append(f"\n## PPT 강의자료\n\n{text}")

    parts.append("\n\n위 내용을 바탕으로 완전한 HTML 강의노트를 생성해주세요.")

    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8192,
        system=[
            {
                "type": "text",
                "text": SYSTEM_PROMPT,
                "cache_control": {"type": "ephemeral"},
            }
        ],
        messages=[{"role": "user", "content": "".join(parts)}],
    )

    html = response.content[0].text.strip()
    # 코드블록으로 감싸진 경우 제거
    html = re.sub(r"^```[a-z]*\n?", "", html)
    html = re.sub(r"\n?```$", "", html)
    return html


# ── 단계별 처리 ──────────────────────────────────────────────

def process_week(subject: str, week_dir: Path) -> bool:
    m = re.match(r"week(\d+)", week_dir.name, re.IGNORECASE)
    if not m:
        print(f"  ⚠️  폴더명 형식 불일치 (weekNN 형식 필요): {week_dir.name}")
        return False

    week_num = int(m.group(1))
    print(f"\n📖 [{subject}] Week {week_num:02d}")

    # 1단계: 입력 수집
    print("  [1/3] 입력 파일 수집...")
    inputs = collect_inputs(week_dir)

    has_content = (
        inputs["transcript"].strip()
        or inputs["pdf"]
        or inputs["pptx"]
    )
    if not has_content:
        print("  ❌ 처리할 파일 없음")
        return False

    file_summary = []
    if inputs["transcript"].strip():
        file_summary.append("녹취스크립트")
    if inputs["pdf"]:
        file_summary.append(f"PDF {len(inputs['pdf'])}개")
    if inputs["pptx"]:
        file_summary.append(f"PPTX {len(inputs['pptx'])}개")
    print(f"       수집: {', '.join(file_summary)}")

    # 2단계: Claude API로 HTML 생성 + 파일 저장
    print("  [2/3] HTML 생성 중 (Claude API)...")
    html = generate_html(subject, week_num, inputs)

    out_folder = OUTPUT_FOLDER_OVERRIDE.get(subject, subject)
    out_dir = ROOT_DIR / "output" / out_folder
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"week{week_num:02d}.html"
    out_path.write_text(html, encoding="utf-8")
    print(f"       저장: {out_path.relative_to(ROOT_DIR)}")

    # 3단계: data.js 업데이트
    print("  [3/3] data.js 업데이트...")
    try:
        _update_data_js(subject, str(week_num))
    except SystemExit as e:
        if e.code != 0:
            print(f"       ⚠️  data.js 업데이트 실패 (알 수 없는 과목 또는 항목 없음)")

    return True


def git_push(processed: list):
    print("\n  [4/4] Git push...")
    label = ", ".join(f"{s} week{w:02d}" for s, w in processed)

    cmds = [
        ["git", "-C", str(ROOT_DIR), "add", "output/", "data.js"],
        ["git", "-C", str(ROOT_DIR), "commit", "-m", f"강의노트 업데이트: {label}"],
        ["git", "-C", str(ROOT_DIR), "push"],
    ]

    for cmd in cmds:
        result = subprocess.run(cmd, capture_output=True, text=True)
        display = " ".join(cmd[3:])
        if result.returncode != 0:
            stderr = result.stderr.strip()
            # "nothing to commit"은 정상
            if "nothing to commit" in stderr or "nothing to commit" in result.stdout:
                print(f"       ℹ️  {display}: 변경 없음")
            else:
                print(f"       ❌ {display} 실패:\n          {stderr}")
        else:
            print(f"       ✅ {display}")


def archive_input(subject: str, week_dir: Path):
    date_str = datetime.now().strftime("%Y%m%d")
    dest = DONE_DIR / subject / f"{week_dir.name}_{date_str}"
    dest.parent.mkdir(parents=True, exist_ok=True)
    shutil.move(str(week_dir), str(dest))
    # 빈 과목 폴더 제거
    parent = week_dir.parent
    if parent.exists() and not any(parent.iterdir()):
        parent.rmdir()


# ── 메인 ─────────────────────────────────────────────────────

def main():
    if not INPUT_DIR.exists():
        print(f"❌ input/ 폴더가 없습니다: {INPUT_DIR}")
        sys.exit(1)

    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("❌ ANTHROPIC_API_KEY 환경변수를 설정해주세요.")
        print("   export ANTHROPIC_API_KEY='sk-ant-...'")
        sys.exit(1)

    # input/ 스캔
    pending = []
    for subject_dir in sorted(INPUT_DIR.iterdir()):
        if not subject_dir.is_dir() or subject_dir.name == "done":
            continue
        subject = subject_dir.name
        if subject not in COURSE_MAP:
            print(f"⚠️  알 수 없는 과목: {subject}")
            print(f"   사용 가능: {', '.join(COURSE_MAP.keys())}")
            continue
        for week_dir in sorted(subject_dir.iterdir()):
            if week_dir.is_dir() and re.match(r"week\d+", week_dir.name, re.IGNORECASE):
                pending.append((subject, week_dir))

    if not pending:
        print("📭 처리할 강의가 없습니다.")
        print(f"   예시: input/기독교철학/week07/transcript.txt")
        sys.exit(0)

    print(f"📚 처리 대상: {len(pending)}개 강의\n")

    processed = []
    for subject, week_dir in pending:
        ok = process_week(subject, week_dir)
        if ok:
            week_num = int(re.search(r"\d+", week_dir.name).group())
            processed.append((subject, week_num))
            archive_input(subject, week_dir)

    if processed:
        git_push(processed)
        print(f"\n🎉 완료: {len(processed)}개 강의노트 업로드됨")
    else:
        print("\n⚠️  처리된 강의 없음")


if __name__ == "__main__":
    main()
